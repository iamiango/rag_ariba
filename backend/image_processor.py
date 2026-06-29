# backend/image_processor.py
"""
图片提取与 OCR 处理模块
从 PDF/PPT/Word 文档中提取嵌入图片并进行 OCR 识别
"""

import os
import uuid
import logging
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime
from io import BytesIO

import numpy as np
from PIL import Image

# 配置日志
logger = logging.getLogger(__name__)

# 检查可用的依赖
try:
    from paddleocr import PaddleOCR
    PADDLE_AVAILABLE = True
except ImportError:
    PADDLE_AVAILABLE = False
    logger.warning("PaddleOCR 未安装，图片 OCR 功能将不可用。安装命令: pip install paddleocr paddlepaddle")

try:
    import fitz  # pymupdf - 用于快速 PDF 图片提取
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger.info("PyMuPDF 未安装，将使用 pdfplumber 提取 PDF 图片")

try:
    import cv2
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False
    logger.warning("OpenCV 未安装，图片预处理功能受限。安装命令: pip install opencv-python-headless")


@dataclass
class ImageExtractionConfig:
    """图片提取与 OCR 配置"""

    # 图片提取设置
    extract_images: bool = True
    min_image_size: int = 50  # 最小像素尺寸，过滤太小的图片
    max_images_per_document: int = 50  # 每个文档最多提取图片数

    # OCR 设置
    enable_ocr: bool = True
    ocr_engine: str = 'paddleocr'  # paddleocr, tesseract
    ocr_languages: List[str] = field(default_factory=lambda: ['ch', 'en'])  # 中英文
    ocr_confidence_threshold: float = 0.3  # OCR 置信度阈值

    # 性能设置
    use_gpu: bool = False  # 是否使用 GPU

    # 输出设置
    save_extracted_images: bool = False  # 是否保存提取的图片
    image_output_dir: str = './extracted_images'
    include_image_context: bool = True  # 是否包含图片周围的文本上下文


class ImageExtractor:
    """从不同文档格式中提取图片"""

    def __init__(self, config: ImageExtractionConfig = None):
        self.config = config or ImageExtractionConfig()

    def extract_from_pdf(self, pdf_path: str) -> List[Dict]:
        """从 PDF 中提取图片"""
        images = []

        # 优先使用 PyMuPDF（更快）
        if PYMUPDF_AVAILABLE:
            images = self._extract_pdf_pymupdf(pdf_path)
        else:
            images = self._extract_pdf_pdfplumber(pdf_path)

        return images

    def _extract_pdf_pymupdf(self, pdf_path: str) -> List[Dict]:
        """使用 PyMuPDF 从 PDF 提取图片（快速）"""
        images = []

        try:
            doc = fitz.open(pdf_path)

            for page_num in range(len(doc)):
                if len(images) >= self.config.max_images_per_document:
                    break

                page = doc[page_num]

                # 获取页面文本作为上下文
                text_context = page.get_text()

                # 获取页面上的所有图片
                image_list = page.get_images()

                for img_idx, img_info in enumerate(image_list):
                    if len(images) >= self.config.max_images_per_document:
                        break

                    try:
                        xref = img_info[0]
                        pix = fitz.Pixmap(doc, xref)

                        # 转换为 RGB（如果是 CMYK）
                        if pix.n - pix.alpha >= 4:
                            pix = fitz.Pixmap(fitz.csRGB, pix)

                        # 检查最小尺寸
                        if pix.width < self.config.min_image_size or pix.height < self.config.min_image_size:
                            continue

                        # 转换为 PIL Image
                        img_data = pix.tobytes("png")
                        pil_image = Image.open(BytesIO(img_data))

                        images.append({
                            'image': pil_image,
                            'page': page_num,
                            'image_index': img_idx,
                            'source': 'pdf',
                            'text_context': text_context[:500] if text_context else '',
                            'document_path': pdf_path,
                            'width': pix.width,
                            'height': pix.height
                        })

                        pix = None  # 释放内存

                    except Exception as e:
                        logger.debug(f"提取 PDF 图片失败 (页 {page_num}, 图 {img_idx}): {str(e)}")

            doc.close()

        except Exception as e:
            logger.error(f"PyMuPDF 提取 PDF 图片失败 {pdf_path}: {str(e)}")

        return images

    def _extract_pdf_pdfplumber(self, pdf_path: str) -> List[Dict]:
        """使用 pdfplumber 从 PDF 提取图片（备用方案）"""
        import pdfplumber

        images = []

        try:
            with pdfplumber.open(pdf_path) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    if len(images) >= self.config.max_images_per_document:
                        break

                    page_images = page.images
                    if not page_images:
                        continue

                    text_context = page.extract_text() or ''

                    for img_idx, img in enumerate(page_images):
                        if len(images) >= self.config.max_images_per_document:
                            break

                        try:
                            # 计算图片尺寸
                            width = img['x1'] - img['x0']
                            height = img['bottom'] - img['top']

                            if width < self.config.min_image_size or height < self.config.min_image_size:
                                continue

                            # 裁剪图片区域
                            cropped = page.crop((img['x0'], img['top'], img['x1'], img['bottom']))
                            pil_image = cropped.to_image().original

                            images.append({
                                'image': pil_image,
                                'page': page_idx,
                                'image_index': img_idx,
                                'source': 'pdf',
                                'text_context': text_context[:500],
                                'document_path': pdf_path,
                                'width': width,
                                'height': height
                            })

                        except Exception as e:
                            logger.debug(f"pdfplumber 提取图片失败: {str(e)}")

        except Exception as e:
            logger.error(f"pdfplumber 提取 PDF 图片失败 {pdf_path}: {str(e)}")

        return images

    def extract_from_pptx(self, pptx_path: str) -> List[Dict]:
        """从 PowerPoint 中提取图片"""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        images = []

        try:
            prs = Presentation(pptx_path)

            for slide_idx, slide in enumerate(prs.slides):
                if len(images) >= self.config.max_images_per_document:
                    break

                # 收集幻灯片文本上下文
                text_parts = []

                for shape in slide.shapes:
                    if len(images) >= self.config.max_images_per_document:
                        break

                    # 收集文本
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                    # 提取图片
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob

                            pil_image = Image.open(BytesIO(image_bytes))

                            # 检查最小尺寸
                            if pil_image.width < self.config.min_image_size or pil_image.height < self.config.min_image_size:
                                continue

                            images.append({
                                'image': pil_image,
                                'slide': slide_idx,
                                'source': 'pptx',
                                'text_context': '\n'.join(text_parts)[:500],
                                'shape_name': shape.name,
                                'document_path': pptx_path,
                                'width': pil_image.width,
                                'height': pil_image.height
                            })

                        except Exception as e:
                            logger.debug(f"提取 PPT 图片失败 (幻灯片 {slide_idx}): {str(e)}")

        except Exception as e:
            logger.error(f"提取 PPTX 图片失败 {pptx_path}: {str(e)}")

        return images

    def extract_from_docx(self, docx_path: str) -> List[Dict]:
        """从 Word 文档中提取图片"""
        from docx import Document

        images = []

        try:
            doc = Document(docx_path)

            # 获取文档文本上下文
            text_context = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])[:500]

            # 从文档关系中提取图片
            for rel_id, rel in doc.part.rels.items():
                if len(images) >= self.config.max_images_per_document:
                    break

                if "image" in rel.target_ref:
                    try:
                        image_bytes = rel.target_part.blob
                        pil_image = Image.open(BytesIO(image_bytes))

                        # 检查最小尺寸
                        if pil_image.width < self.config.min_image_size or pil_image.height < self.config.min_image_size:
                            continue

                        images.append({
                            'image': pil_image,
                            'source': 'docx',
                            'text_context': text_context,
                            'relation_id': rel_id,
                            'document_path': docx_path,
                            'width': pil_image.width,
                            'height': pil_image.height
                        })

                    except Exception as e:
                        logger.debug(f"提取 DOCX 图片失败: {str(e)}")

        except Exception as e:
            logger.error(f"提取 DOCX 图片失败 {docx_path}: {str(e)}")

        return images

    def extract_images(self, document_path: str) -> List[Dict]:
        """根据文件类型提取图片"""
        file_ext = Path(document_path).suffix.lower()

        if file_ext == '.pdf':
            return self.extract_from_pdf(document_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self.extract_from_pptx(document_path)
        elif file_ext in ['.docx', '.doc']:
            return self.extract_from_docx(document_path)
        else:
            logger.warning(f"不支持的文件格式: {file_ext}")
            return []


class OCRProcessor:
    """OCR 处理器"""

    _instance = None  # 单例模式，避免重复加载模型

    def __new__(cls, config: ImageExtractionConfig = None):
        if cls._instance is None:
            cls._instance = super().__new__(cls)
            cls._instance._initialized = False
        return cls._instance

    def __init__(self, config: ImageExtractionConfig = None):
        if self._initialized:
            return

        self.config = config or ImageExtractionConfig()
        self.ocr_engine = None
        self._initialize_ocr_engine()
        self._initialized = True

    def _initialize_ocr_engine(self):
        """初始化 OCR 引擎"""
        if not self.config.enable_ocr:
            logger.info("OCR 已禁用")
            return

        if self.config.ocr_engine == 'paddleocr' and PADDLE_AVAILABLE:
            try:
                # PaddleOCR 3.x 新版 API（简化参数）
                self.ocr_engine = PaddleOCR(lang='ch')  # 中文模型同时支持英文
                logger.info("PaddleOCR 初始化成功")
            except Exception as e:
                logger.error(f"PaddleOCR 初始化失败: {str(e)}")
                self.ocr_engine = None
        else:
            logger.warning(f"OCR 引擎 {self.config.ocr_engine} 不可用")

    def preprocess_image(self, image: Image.Image) -> Image.Image:
        """预处理图片以提高 OCR 准确率"""
        if not CV2_AVAILABLE:
            return image

        try:
            # 转换为 numpy 数组
            image_np = np.array(image)

            # 转换为灰度图（如果是彩色）
            if len(image_np.shape) == 3:
                gray = cv2.cvtColor(image_np, cv2.COLOR_RGB2GRAY)
            else:
                gray = image_np

            # 增强对比度 (CLAHE)
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            enhanced = clahe.apply(gray)

            # 降噪
            denoised = cv2.fastNlMeansDenoising(enhanced)

            return Image.fromarray(denoised)

        except Exception as e:
            logger.debug(f"图片预处理失败: {str(e)}")
            return image

    def process_image(self, image: Image.Image) -> Dict:
        """对图片进行 OCR 识别"""
        if not self.config.enable_ocr or self.ocr_engine is None:
            return {'text': '', 'confidence': 0, 'error': 'OCR 未启用或初始化失败'}

        try:
            # 转换为 numpy 数组
            image_np = np.array(image)

            # 确保是 RGB 格式
            if len(image_np.shape) == 2:
                image_np = cv2.cvtColor(image_np, cv2.COLOR_GRAY2RGB) if CV2_AVAILABLE else np.stack([image_np]*3, axis=-1)
            elif image_np.shape[2] == 4:
                image_np = image_np[:, :, :3]

            # 执行 OCR (PaddleOCR 3.x 新版 API)
            result = self.ocr_engine.ocr(image_np)

            # 解析结果 - 新版 PaddleOCR 3.x 返回格式
            text_lines = []
            confidences = []

            if result and len(result) > 0:
                # 新版返回格式: result 是一个列表，每个元素是一个字典
                for item in result:
                    if isinstance(item, dict):
                        # 新版格式: {'rec_texts': [...], 'rec_scores': [...]}
                        rec_texts = item.get('rec_texts', [])
                        rec_scores = item.get('rec_scores', [])

                        for text, score in zip(rec_texts, rec_scores):
                            if score > self.config.ocr_confidence_threshold:
                                text_lines.append(text)
                                confidences.append(score)
                    elif isinstance(item, list):
                        # 旧版格式兼容: [[box, (text, score)], ...]
                        for line in item:
                            if line and len(line) >= 2:
                                text = line[1][0] if isinstance(line[1], tuple) else line[1]
                                confidence = line[1][1] if isinstance(line[1], tuple) else 0.9

                                if confidence > self.config.ocr_confidence_threshold:
                                    text_lines.append(text)
                                    confidences.append(confidence)

            return {
                'text': ' '.join(text_lines),
                'confidence': float(np.mean(confidences)) if confidences else 0,
                'line_count': len(text_lines)
            }

        except Exception as e:
            logger.error(f"OCR 处理失败: {str(e)}")
            return {'text': '', 'confidence': 0, 'error': str(e)}


class ImageDocumentLoader:
    """集成图片提取 + OCR + LangChain Document 创建"""

    def __init__(self, config: ImageExtractionConfig = None):
        self.config = config or ImageExtractionConfig()
        self.extractor = ImageExtractor(self.config)
        self.ocr_processor = OCRProcessor(self.config)

        # 创建输出目录
        if self.config.save_extracted_images:
            os.makedirs(self.config.image_output_dir, exist_ok=True)

    def save_image(self, image: Image.Image, source: str, index: int, doc_name: str) -> Optional[str]:
        """保存提取的图片"""
        if not self.config.save_extracted_images:
            return None

        try:
            filename = f"{doc_name}_{source}_{index}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
            filepath = os.path.join(self.config.image_output_dir, filename)
            image.save(filepath, format='PNG')
            return filepath
        except Exception as e:
            logger.debug(f"保存图片失败: {str(e)}")
            return None

    def load_images_from_document(self, document_path: str) -> List:
        """从文档中提取图片并进行 OCR，返回 LangChain Document 列表"""
        from langchain.schema import Document

        # 提取图片
        extracted_images = self.extractor.extract_images(document_path)

        if not extracted_images:
            return []

        documents = []
        doc_name = Path(document_path).stem

        for idx, image_data in enumerate(extracted_images):
            image = image_data['image']

            # 保存图片（如果配置了）
            image_path = self.save_image(image, image_data['source'], idx, doc_name)

            # OCR 识别
            ocr_result = self.ocr_processor.process_image(image)

            # 跳过没有识别出文字的图片
            if not ocr_result['text'].strip():
                continue

            # 构建文档内容
            content_parts = []

            # 添加图片上下文
            if self.config.include_image_context and image_data.get('text_context'):
                content_parts.append(f"[图片上下文] {image_data['text_context']}")

            # 添加 OCR 文本
            content_parts.append(f"[图片内容] {ocr_result['text']}")

            # 创建 LangChain Document
            page_or_slide = image_data.get('page', image_data.get('slide', 0))

            doc = Document(
                page_content='\n'.join(content_parts),
                metadata={
                    'source': document_path,
                    'type': 'image_ocr',
                    'image_source': image_data['source'],
                    'image_path': image_path,
                    'image_index': idx,
                    'page': page_or_slide,
                    'ocr_confidence': ocr_result['confidence'],
                    'image_width': image_data.get('width', 0),
                    'image_height': image_data.get('height', 0),
                    'extraction_time': datetime.now().isoformat()
                }
            )

            documents.append(doc)

        if documents:
            logger.info(f"从 {document_path} 提取并识别了 {len(documents)} 张图片")

        return documents


# 便捷函数
def extract_images_with_ocr(document_path: str, config: ImageExtractionConfig = None) -> List:
    """
    便捷函数：从文档中提取图片并进行 OCR

    Args:
        document_path: 文档路径
        config: 配置对象（可选）

    Returns:
        LangChain Document 列表
    """
    loader = ImageDocumentLoader(config)
    return loader.load_images_from_document(document_path)


def is_ocr_available() -> bool:
    """检查 OCR 功能是否可用"""
    return PADDLE_AVAILABLE


def get_supported_formats() -> List[str]:
    """获取支持图片提取的文件格式"""
    return ['.pdf', '.pptx', '.ppt', '.docx', '.doc']
